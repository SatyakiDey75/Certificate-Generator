from pptx import Presentation
import os
import csv
def edit_powerpoint_text(pptx_file_path, old_text, new_text):
    presentation = Presentation(pptx_file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
    output_pptx = os.path.join("output", f"{new_text}.pptx")
    presentation.save(output_pptx)

if __name__ == "__main__":
    
    os.makedirs("output", exist_ok=True)

    # change this with the path of the ppt template
    pptx_file_path = "template.pptx"

    # change this with the path of the file having student names
    student_names_csv = "file_path.csv"

    # change this with the text that is to be replaced with the student's name
    text_to_be_replaced = "text_to_be_replaced"
    
    with open(student_names_csv, 'r') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                edit_powerpoint_text(pptx_file_path, text_to_be_replaced, row[0])


# After this run the following commands in the terminal:
# cd "Certificate Generator"
# ppt2pdf dir output

